import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from io import BytesIO
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from htmlTemplates import css, bot_template, user_template

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def get_docx_text(docx_docs):
    text = ""
    for docx_file in docx_docs:
        doc = DocxDocument(docx_file)
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
    return text

def get_pptx_text(pptx_docs):
    text = ""
    for pptx_file in pptx_docs:
        presentation = Presentation(pptx_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_image_text(image_docs):
    text = ""
    for image_file in image_docs:
        image = Image.open(image_file)
        buffer = BytesIO()
        image.save(buffer, format="JPEG")
        text += f"Image content: {image_file.name}\n"
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    load_dotenv()
    st.set_page_config(page_title="Text-Tract", page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("Text-Tract :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        file_types = ["pdf", "docx", "pptx", "jpg", "jpeg", "png"]
        files = st.file_uploader(
            "Upload your files here and click on 'Process'", accept_multiple_files=True, type=file_types)

        if st.button("Process"):
            with st.spinner("Processing"):
                raw_text = ""

                # Process uploaded files based on their type
                pdf_docs = [f for f in files if f.type == "application/pdf"]
                docx_docs = [f for f in files if f.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]
                pptx_docs = [f for f in files if f.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
                image_docs = [f for f in files if f.type in ["image/jpeg", "image/png"]]

                if pdf_docs:
                    raw_text += get_pdf_text(pdf_docs)
                if docx_docs:
                    raw_text += get_docx_text(docx_docs)
                if pptx_docs:
                    raw_text += get_pptx_text(pptx_docs)
                if image_docs:
                    raw_text += get_image_text(image_docs)

                # get the text chunks
                text_chunks = get_text_chunks(raw_text)

                # create vector store
                vectorstore = get_vectorstore(text_chunks)

                # create conversation chain
                st.session_state.conversation = get_conversation_chain(vectorstore)

if __name__ == '__main__':
    main()

